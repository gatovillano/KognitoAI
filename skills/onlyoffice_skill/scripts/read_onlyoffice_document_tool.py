# skills/onlyoffice_skill/scripts/read_onlyoffice_document_tool.py

import os
import logging
import uuid
import csv
from io import BytesIO
from typing import Any, Type, Optional
from pydantic import BaseModel, Field
from langchain_core.tools import BaseTool
from sqlalchemy import select
from core.database import SessionLocal, Document
from utils.document_parser import extract_text_and_metadata_from_document

try:
    from openpyxl import load_workbook
except ImportError:
    load_workbook = None

try:
    from pptx import Presentation
except ImportError:
    Presentation = None

logger = logging.getLogger(__name__)

DEFAULT_DOCS_ROOT = os.path.join("/app/media", "documents")
DOCUMENTS_ROOT = os.environ.get("ONLYOFFICE_DOCS_ROOT", DEFAULT_DOCS_ROOT)


class ReadOnlyOfficeInput(BaseModel):
    document_id: str = Field(..., description="ID UUID del documento que se desea leer.")


class ReadOnlyOfficeDocumentTool(BaseTool):
    name: str = "read_onlyoffice_document"
    description: str = (
        "Lee el contenido textual completo de un documento almacenado en el servidor. "
        "Soporta PDF (incluyendo escaneados con OCR), Word (.docx), Excel (.xlsx), "
        "PowerPoint (.pptx), CSV y TXT. Úsala cuando tengas el ID UUID del documento."
    )
    args_schema: Type[BaseModel] = ReadOnlyOfficeInput
    account_id: str = Field(..., description="ID de cuenta del usuario, inyectado automáticamente.")

    async def _arun(self, document_id: str, **kwargs: Any) -> str:
        try:
            acc_id = uuid.UUID(self.account_id)
            doc_id = uuid.UUID(document_id)

            async with SessionLocal() as db:
                stmt = select(Document).where(
                    Document.id == doc_id,
                    Document.account_id == acc_id
                )
                result = await db.execute(stmt)
                doc = result.scalar_one_or_none()

            if not doc:
                return f"Documento con ID '{document_id}' no encontrado o no tienes permiso para leerlo."

            # Resolver ruta física
            file_path = doc.file_path
            if not os.path.isabs(file_path):
                file_path = os.path.join(DOCUMENTS_ROOT, file_path)

            if not os.path.exists(file_path):
                # Fallback: buscar relativo al cwd
                for candidate in [
                    os.path.join(os.getcwd(), "media", "documents", doc.file_path),
                    os.path.join(os.getcwd(), doc.file_path),
                ]:
                    if os.path.exists(candidate):
                        file_path = candidate
                        break
                else:
                    return (
                        f"El archivo físico '{doc.filename}' no se encuentra en el servidor. "
                        f"Ruta buscada: {file_path}"
                    )

            with open(file_path, "rb") as f:
                file_bytes = f.read()

            ext = doc.extension.lower().replace(".", "")
            header = f"--- CONTENIDO DE '{doc.filename}' ---\n\n"

            # --- Formatos manejados por el parser central (PDF, DOCX, TXT, MD, imágenes) ---
            PARSER_EXTENSIONS = {"pdf", "docx", "txt", "md", "png", "jpg", "jpeg", "webp"}
            if ext in PARSER_EXTENSIONS:
                text, _ = await extract_text_and_metadata_from_document(doc.filename, file_bytes)
                if text:
                    return header + text
                return header + f"⚠️ No se pudo extraer texto del archivo '{doc.filename}'."

            # --- Excel (.xlsx) ---
            if ext == "xlsx":
                if not load_workbook:
                    return header + "⚠️ La librería 'openpyxl' no está instalada."
                try:
                    wb = load_workbook(BytesIO(file_bytes), data_only=True)
                    parts = []
                    for sheet_name in wb.sheetnames:
                        ws = wb[sheet_name]
                        parts.append(f"[Hoja: {sheet_name}]")
                        for row in ws.iter_rows(min_row=1, values_only=True):
                            if any(cell is not None for cell in row):
                                parts.append(" | ".join("" if c is None else str(c) for c in row))
                    return header + "\n".join(parts)
                except Exception as e:
                    return header + f"❌ Error al procesar XLSX: {e}"

            # --- PowerPoint (.pptx) ---
            if ext == "pptx":
                if not Presentation:
                    return header + "⚠️ La librería 'python-pptx' no está instalada."
                try:
                    prs = Presentation(BytesIO(file_bytes))
                    parts = []
                    for i, slide in enumerate(prs.slides, 1):
                        slide_texts = [
                            shape.text for shape in slide.shapes
                            if hasattr(shape, "text") and shape.text.strip()
                        ]
                        if slide_texts:
                            parts.append(f"[Slide {i}]")
                            parts.extend(slide_texts)
                    return header + "\n".join(parts)
                except Exception as e:
                    return header + f"❌ Error al procesar PPTX: {e}"

            # --- CSV ---
            if ext == "csv":
                try:
                    text = file_bytes.decode("utf-8", errors="replace")
                    reader = csv.reader(text.splitlines())
                    rows = [", ".join(row) for row in reader]
                    return header + "\n".join(rows)
                except Exception as e:
                    return header + f"❌ Error al leer CSV: {e}"

            # --- Fallback genérico texto plano ---
            try:
                # Rechazar binarios evidentes (ZIP = Office moderno, OLE2 = Office antiguo)
                if file_bytes[:4] in (b'PK\x03\x04', b'\xd0\xcf\x11\xe0'):
                    return (
                        header +
                        f"⚠️ El archivo '{doc.filename}' ({ext}) es un formato binario no soportado directamente. "
                        "Considera convertirlo a PDF o DOCX antes de subirlo."
                    )
                decoded = file_bytes.decode("utf-8", errors="replace")
                return header + decoded
            except Exception as e:
                return header + f"❌ Error al leer el archivo: {e}"

        except Exception as e:
            logger.error(f"Error en ReadOnlyOfficeDocumentTool para doc_id='{document_id}': {e}", exc_info=True)
            return f"Error inesperado al leer el documento: {e}"

    def _run(self, *args: Any, **kwargs: Any) -> Any:
        raise NotImplementedError("Esta herramienta solo soporta ejecución asíncrona.")
